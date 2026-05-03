"""Build two editable PPTX sales decks: Individuals and Organizations.

Output:
- deliverables/sales_decks/Individuals_Deck.pptx
- deliverables/sales_decks/Organizations_Deck.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image
import os
import tempfile

ASSETS_DIR = "attached_assets"
_BUILD_TMP = os.path.join(tempfile.gettempdir(), "sales_deck_screenshots")
os.makedirs(_BUILD_TMP, exist_ok=True)

# ---------- Brand system ----------
BG_DEEP    = RGBColor(0x0A, 0x0F, 0x1C)   # deepest background
BG         = RGBColor(0x0F, 0x17, 0x2A)   # base slide bg
SURFACE    = RGBColor(0x16, 0x21, 0x3A)   # card surface
SURFACE_2  = RGBColor(0x1E, 0x2C, 0x4A)   # elevated card
LINE       = RGBColor(0x2A, 0x3A, 0x5C)   # subtle border
PRIMARY    = RGBColor(0x1F, 0xA8, 0xE0)   # cyan accent
PRIMARY_DK = RGBColor(0x0E, 0x6E, 0x96)
ACCENT     = RGBColor(0xF2, 0x6D, 0x55)   # coral
SUCCESS    = RGBColor(0x22, 0xC5, 0x5E)
WARN       = RGBColor(0xF5, 0xB7, 0x2E)
DANGER     = RGBColor(0xEF, 0x55, 0x55)
TEXT       = RGBColor(0xE7, 0xEC, 0xF5)
TEXT_MUTED = RGBColor(0x9A, 0xA6, 0xC0)
TEXT_DIM   = RGBColor(0x6A, 0x78, 0x95)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Inter"
FONT_FALLBACK = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---------- Helpers ----------
def _set_fill(shape, color, transparency=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def _no_line(shape):
    shape.line.fill.background()

def _line(shape, color, width=0.75):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)

def add_rect(slide, x, y, w, h, fill=SURFACE, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        _set_fill(s, fill)
    if line is None:
        _no_line(s)
    else:
        _line(s, line, line_w)
    return s

def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=14, color=TEXT, bullet_color=PRIMARY,
                line_spacing=1.35, gap_pt=6):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if i > 0:
            p.space_before = Pt(gap_pt)
        r1 = p.add_run()
        r1.text = "▸  "
        r1.font.name = FONT
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color
        r2 = p.add_run()
        r2.text = item
        r2.font.name = FONT
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return tb

# ---------- Slide chrome ----------
def base_slide(prs, *, bg=BG, accent_strip=True):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    # background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=bg)
    # subtle deeper band along top
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), fill=PRIMARY)
    return slide

def add_footer(slide, page_num, total, label="Roster Intelligence Platform"):
    add_text(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
             label, size=9, color=TEXT_DIM)
    add_text(slide, Inches(11.5), Inches(7.1), Inches(1.4), Inches(0.3),
             f"{page_num:02d} / {total:02d}", size=9, color=TEXT_DIM, align=PP_ALIGN.RIGHT)

def add_eyebrow(slide, x, y, text, color=PRIMARY):
    # tiny accent line + uppercase label
    add_rect(slide, x, y + Inches(0.08), Inches(0.28), Inches(0.04), fill=color)
    add_text(slide, x + Inches(0.42), y, Inches(6), Inches(0.3),
             text.upper(), size=10, bold=True, color=color)

def add_slide_header(slide, eyebrow, title, subtitle=None):
    add_eyebrow(slide, Inches(0.6), Inches(0.55), eyebrow)
    add_text(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.9),
             title, size=32, bold=True, color=TEXT, line_spacing=1.05)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.85), Inches(12), Inches(0.5),
                 subtitle, size=14, color=TEXT_MUTED)
    # underline accent
    add_rect(slide, Inches(0.6), Inches(2.45), Inches(0.6), Inches(0.05), fill=ACCENT)

# ---------- Specialized slides ----------
def cover_slide(prs, eyebrow, title, subtitle, edition_label):
    s = base_slide(prs, bg=BG_DEEP)
    # left accent column
    add_rect(s, 0, 0, Inches(0.35), SLIDE_H, fill=PRIMARY)
    # decorative diagonal blocks (tactical look)
    add_rect(s, Inches(11.3), Inches(0.7), Inches(1.6), Inches(0.06), fill=PRIMARY)
    add_rect(s, Inches(11.3), Inches(0.86), Inches(0.9), Inches(0.04), fill=ACCENT)
    add_rect(s, Inches(11.3), Inches(0.96), Inches(0.4), Inches(0.04), fill=TEXT_DIM)
    # eyebrow
    add_eyebrow(s, Inches(1.0), Inches(2.4), eyebrow)
    # main title
    add_text(s, Inches(1.0), Inches(2.85), Inches(11), Inches(2.2),
             title, size=58, bold=True, color=WHITE, line_spacing=1.0)
    # subtitle
    add_text(s, Inches(1.0), Inches(5.05), Inches(10.5), Inches(1.0),
             subtitle, size=18, color=TEXT_MUTED, line_spacing=1.3)
    # edition tag chip
    chip = add_rect(s, Inches(1.0), Inches(6.2), Inches(3.6), Inches(0.45),
                    fill=None, line=PRIMARY, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(1.0), Inches(6.22), Inches(3.6), Inches(0.42),
             edition_label.upper(), size=10, bold=True, color=PRIMARY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # footer line
    add_text(s, Inches(1.0), Inches(6.85), Inches(11), Inches(0.3),
             "Confidential sales material   ·   Built for esports orgs, coaches & competitive players",
             size=10, color=TEXT_DIM)

def section_divider(prs, number, label, headline):
    s = base_slide(prs, bg=BG_DEEP)
    add_text(s, Inches(0.7), Inches(2.2), Inches(2), Inches(0.6),
             f"§ {number:02d}", size=22, bold=True, color=PRIMARY)
    add_rect(s, Inches(0.7), Inches(2.95), Inches(0.7), Inches(0.05), fill=ACCENT)
    add_text(s, Inches(0.7), Inches(3.15), Inches(2.4), Inches(0.4),
             label.upper(), size=12, bold=True, color=TEXT_MUTED)
    add_text(s, Inches(0.7), Inches(3.65), Inches(12), Inches(2.5),
             headline, size=44, bold=True, color=WHITE, line_spacing=1.05)

def feature_card(slide, x, y, w, h, icon, title, body):
    add_rect(slide, x, y, w, h, fill=SURFACE, line=LINE)
    # icon badge
    add_rect(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.55), Inches(0.55),
             fill=PRIMARY_DK, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.55), Inches(0.55),
             icon, size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(0.3), y + Inches(1.0), w - Inches(0.6), Inches(0.4),
             title, size=14, bold=True, color=TEXT)
    add_text(slide, x + Inches(0.3), y + Inches(1.45), w - Inches(0.6), h - Inches(1.6),
             body, size=11, color=TEXT_MUTED, line_spacing=1.35)

def stat_card(slide, x, y, w, h, value, label, accent=PRIMARY):
    add_rect(slide, x, y, w, h, fill=SURFACE, line=LINE)
    add_rect(slide, x, y, Inches(0.08), h, fill=accent)
    add_text(slide, x + Inches(0.3), y + Inches(0.3), w - Inches(0.6), Inches(0.9),
             value, size=34, bold=True, color=WHITE)
    add_text(slide, x + Inches(0.3), y + Inches(1.25), w - Inches(0.6), h - Inches(1.4),
             label, size=11, color=TEXT_MUTED, line_spacing=1.3)

def comparison_table(slide, x, y, w, h, headers, rows, highlight_col=2):
    """headers: list[str], rows: list[list[str]]; col 0 is row label, others are columns."""
    n_cols = len(headers)
    n_rows = len(rows) + 1
    col_w = w / n_cols
    row_h = h / n_rows
    # header row
    for i, hd in enumerate(headers):
        fill = PRIMARY if i == highlight_col else SURFACE_2
        color = WHITE if i == highlight_col else TEXT
        add_rect(slide, x + col_w * i, y, col_w, row_h, fill=fill, line=LINE)
        add_text(slide, x + col_w * i + Inches(0.15), y, col_w - Inches(0.3), row_h,
                 hd, size=11, bold=True, color=color,
                 align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.MIDDLE)
    # body rows
    for r, row in enumerate(rows):
        ry = y + row_h * (r + 1)
        for i, cell in enumerate(row):
            row_fill = SURFACE if r % 2 == 0 else BG
            if i == highlight_col:
                row_fill = SURFACE_2
            add_rect(slide, x + col_w * i, ry, col_w, row_h, fill=row_fill, line=LINE)
            color = WHITE if i == highlight_col else (TEXT if i == 0 else TEXT_MUTED)
            bold = (i == 0)
            add_text(slide, x + col_w * i + Inches(0.15), ry, col_w - Inches(0.3), row_h,
                     cell, size=10, bold=bold, color=color,
                     align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT,
                     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)

def cta_slide(prs, eyebrow, headline, subhead, primary_label, secondary_label, footnote):
    s = base_slide(prs, bg=BG_DEEP)
    add_eyebrow(s, Inches(0.7), Inches(2.3), eyebrow, color=ACCENT)
    add_text(s, Inches(0.7), Inches(2.75), Inches(12), Inches(2.2),
             headline, size=52, bold=True, color=WHITE, line_spacing=1.05)
    add_text(s, Inches(0.7), Inches(4.95), Inches(11), Inches(0.9),
             subhead, size=16, color=TEXT_MUTED, line_spacing=1.35)
    # CTA buttons
    btn_y = Inches(6.0)
    # primary
    add_rect(s, Inches(0.7), btn_y, Inches(3.6), Inches(0.65),
             fill=PRIMARY, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(0.7), btn_y, Inches(3.6), Inches(0.65),
             primary_label, size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # secondary
    add_rect(s, Inches(4.5), btn_y, Inches(3.6), Inches(0.65),
             fill=None, line=PRIMARY, line_w=1.25, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(4.5), btn_y, Inches(3.6), Inches(0.65),
             secondary_label, size=13, bold=True, color=PRIMARY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.3),
             footnote, size=10, color=TEXT_DIM)
    return s

def closing_slide(prs, headline, line2, contact):
    s = base_slide(prs, bg=BG_DEEP)
    add_rect(s, 0, 0, Inches(0.35), SLIDE_H, fill=ACCENT)
    add_text(s, Inches(1.0), Inches(2.7), Inches(11), Inches(1.6),
             headline, size=64, bold=True, color=WHITE, line_spacing=1.0)
    add_text(s, Inches(1.0), Inches(4.5), Inches(11), Inches(0.6),
             line2, size=18, color=TEXT_MUTED)
    add_rect(s, Inches(1.0), Inches(5.4), Inches(0.7), Inches(0.05), fill=PRIMARY)
    add_text(s, Inches(1.0), Inches(5.55), Inches(11), Inches(0.5),
             contact, size=13, bold=True, color=PRIMARY)

def _prepare_screenshot(image_path, target_w_emu, target_h_emu, key):
    """Center-crop the given image to match the target frame aspect ratio,
    preserving as much of the original product UI as possible. Returns the
    path to a cached PNG ready to embed in the PPTX."""
    target_ratio = target_w_emu / target_h_emu
    cache_path = os.path.join(_BUILD_TMP, f"{key}.png")
    if os.path.exists(cache_path):
        return cache_path
    im = Image.open(image_path).convert("RGB")
    iw, ih = im.size
    src_ratio = iw / ih
    if src_ratio > target_ratio:
        # source is wider — crop sides
        new_w = int(ih * target_ratio)
        x0 = (iw - new_w) // 2
        im = im.crop((x0, 0, x0 + new_w, ih))
    elif src_ratio < target_ratio:
        # source is taller — crop top/bottom, biased toward top so headers stay
        new_h = int(iw / target_ratio)
        y0 = max(0, (ih - new_h) // 4)  # bias upward
        im = im.crop((0, y0, iw, y0 + new_h))
    # downscale very large images to keep file size reasonable
    max_w = 2200
    if im.size[0] > max_w:
        ratio = max_w / im.size[0]
        im = im.resize((max_w, int(im.size[1] * ratio)), Image.LANCZOS)
    im.save(cache_path, format="PNG", optimize=True)
    return cache_path


def mockup_frame(slide, x, y, w, h, caption=None, image=None):
    """Draw a dark window-chrome frame around a product screenshot.

    If `image` is provided, a real screenshot from `attached_assets/` is
    embedded inside the frame body. Otherwise the function falls back to a
    schematic placeholder grid so the deck still builds without assets.
    """
    # outer frame + title bar (kept for both modes so styling is identical)
    add_rect(slide, x, y, w, h, fill=SURFACE, line=LINE)
    add_rect(slide, x, y, w, Inches(0.32), fill=SURFACE_2, line=LINE)
    add_rect(slide, x + Inches(0.15), y + Inches(0.10), Inches(0.13), Inches(0.13),
             fill=DANGER, shape=MSO_SHAPE.OVAL)
    add_rect(slide, x + Inches(0.32), y + Inches(0.10), Inches(0.13), Inches(0.13),
             fill=WARN, shape=MSO_SHAPE.OVAL)
    add_rect(slide, x + Inches(0.49), y + Inches(0.10), Inches(0.13), Inches(0.13),
             fill=SUCCESS, shape=MSO_SHAPE.OVAL)

    body_y = y + Inches(0.32)
    body_h = h - Inches(0.32)

    if image:
        image_path = image if os.path.isabs(image) else os.path.join(ASSETS_DIR, image)
        if not os.path.exists(image_path):
            raise FileNotFoundError(f"mockup_frame: screenshot not found: {image_path}")
        # small inset so the screenshot doesn't touch the chrome
        inset = Inches(0.06)
        ix, iy = x + inset, body_y + inset
        iw_emu, ih_emu = w - inset * 2, body_h - inset * 2
        # include source mtime in cache key so updated screenshots invalidate cleanly
        mtime = int(os.path.getmtime(image_path))
        key = (os.path.splitext(os.path.basename(image_path))[0]
               + f"_{int(iw_emu)}x{int(ih_emu)}_{mtime}")
        cropped = _prepare_screenshot(image_path, int(iw_emu), int(ih_emu), key)
        # background behind the image, in case PNG has transparency
        add_rect(slide, ix, iy, iw_emu, ih_emu, fill=BG_DEEP, line=None)
        slide.shapes.add_picture(cropped, ix, iy, width=iw_emu, height=ih_emu)
    else:
        # schematic placeholder (legacy)
        nav_w = Inches(1.1)
        add_rect(slide, x, body_y, nav_w, body_h, fill=BG_DEEP, line=LINE)
        for i in range(5):
            add_rect(slide, x + Inches(0.15), y + Inches(0.55) + Inches(0.45) * i,
                     nav_w - Inches(0.3), Inches(0.22), fill=SURFACE_2)
        cx = x + nav_w + Inches(0.25)
        cw = w - nav_w - Inches(0.4)
        add_rect(slide, cx, y + Inches(0.5), cw, Inches(0.5), fill=SURFACE_2)
        card_h = (h - Inches(1.6)) / 2
        add_rect(slide, cx, y + Inches(1.15), cw / 2 - Inches(0.1), card_h, fill=SURFACE_2)
        add_rect(slide, cx + cw / 2 + Inches(0.1), y + Inches(1.15),
                 cw / 2 - Inches(0.1), card_h, fill=SURFACE_2)
        add_rect(slide, cx, y + Inches(1.3) + card_h, cw, h - Inches(1.5) - card_h, fill=SURFACE_2)

    if caption:
        add_text(slide, x, y + h + Inches(0.1), w, Inches(0.3),
                 caption, size=10, color=TEXT_DIM, align=PP_ALIGN.CENTER)

def diagram_pipeline(slide, x, y, w, h, steps):
    n = len(steps)
    box_w = (w - Inches(0.3) * (n - 1)) / n
    for i, (icon, title, body) in enumerate(steps):
        bx = x + (box_w + Inches(0.3)) * i
        add_rect(slide, bx, y, box_w, h, fill=SURFACE, line=LINE)
        add_rect(slide, bx, y, box_w, Inches(0.08), fill=PRIMARY if i % 2 == 0 else ACCENT)
        add_text(slide, bx + Inches(0.25), y + Inches(0.25), box_w - Inches(0.5), Inches(0.4),
                 f"STEP {i+1:02d}", size=9, bold=True, color=PRIMARY)
        add_text(slide, bx + Inches(0.25), y + Inches(0.65), box_w - Inches(0.5), Inches(0.5),
                 title, size=14, bold=True, color=TEXT)
        add_text(slide, bx + Inches(0.25), y + Inches(1.2), box_w - Inches(0.5), h - Inches(1.4),
                 body, size=10, color=TEXT_MUTED, line_spacing=1.3)
        if i < n - 1:
            ax = x + (box_w + Inches(0.3)) * (i + 1) - Inches(0.28)
            add_rect(slide, ax, y + h / 2 - Inches(0.025), Inches(0.26), Inches(0.05),
                     fill=PRIMARY)

def bar_chart(slide, x, y, w, h, title, bars, max_val=None):
    """Simple visual bar chart using shapes (editable)."""
    add_rect(slide, x, y, w, h, fill=SURFACE, line=LINE)
    add_text(slide, x + Inches(0.3), y + Inches(0.25), w - Inches(0.6), Inches(0.4),
             title, size=12, bold=True, color=TEXT)
    inner_x = x + Inches(0.3)
    inner_y = y + Inches(0.85)
    inner_w = w - Inches(0.6)
    inner_h = h - Inches(1.1)
    if max_val is None:
        max_val = max(v for _, v, _ in bars)
    n = len(bars)
    row_h = inner_h / n
    label_w = Inches(2.2)
    track_x = inner_x + label_w
    track_w = inner_w - label_w - Inches(0.7)
    for i, (label, val, color) in enumerate(bars):
        ry = inner_y + row_h * i + Inches(0.08)
        rh = row_h - Inches(0.16)
        add_text(slide, inner_x, ry, label_w - Inches(0.1), rh,
                 label, size=10, color=TEXT_MUTED, anchor=MSO_ANCHOR.MIDDLE)
        # track
        add_rect(slide, track_x, ry + rh / 2 - Inches(0.09), track_w, Inches(0.18),
                 fill=BG_DEEP, line=None)
        # bar
        bw = Emu(int(track_w * (val / max_val)))
        add_rect(slide, track_x, ry + rh / 2 - Inches(0.09), bw, Inches(0.18),
                 fill=color, line=None)
        # value
        add_text(slide, track_x + track_w + Inches(0.1), ry,
                 Inches(0.6), rh, f"{val}", size=10, bold=True, color=TEXT,
                 anchor=MSO_ANCHOR.MIDDLE)

# =====================================================================
#                       INDIVIDUALS DECK (16 slides)
# =====================================================================
def build_individuals(out_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    TOTAL = 16

    # 01 Cover
    cover_slide(
        prs,
        eyebrow="Roster Intelligence  ·  Single Roster Edition",
        title="Stop guessing.\nStart winning your matches.",
        subtitle="A professional-grade roster, stats and prep workspace —\nbuilt for the player who takes the game seriously.",
        edition_label="Individuals  ·  Single Roster",
    )

    # 02 What this is
    s = base_slide(prs)
    add_slide_header(s, "01  ·  What this is",
                     "Your personal command center for one competitive roster.",
                     "Everything you need to track performance, prepare for opponents and improve — in one place.")
    feature_card(s, Inches(0.6),  Inches(3.0), Inches(4.0), Inches(3.6),
                 "01", "One Roster, Fully Owned",
                 "Your players, your matches, your stats, your scouting notes. Private to you, organized like a pro team would.")
    feature_card(s, Inches(4.85), Inches(3.0), Inches(4.0), Inches(3.6),
                 "02", "Real Match Data",
                 "Log results, drafts, maps, heroes/agents and KDA — and watch trends emerge automatically as you play.")
    feature_card(s, Inches(9.10), Inches(3.0), Inches(3.7), Inches(3.6),
                 "03", "Pre-Match Prep",
                 "Scout opponents, study past meetups, and walk into your next match knowing exactly what works.")
    add_footer(s, 2, TOTAL)

    # 03 Why one roster matters
    s = base_slide(prs)
    add_slide_header(s, "02  ·  Why one roster matters",
                     "Most players lose matches before they start.",
                     "Memory, group chats and screenshots don't scale past a few weeks of competitive play.")
    add_bullets(s, Inches(0.6), Inches(3.0), Inches(6.0), Inches(4.0), [
        "You forget which comps actually work for your line-up",
        "You can't see your real win-rate by map, mode or opponent",
        "Scouting an opponent means scrolling through Discord history",
        "Improvement feels invisible — there's nothing to point at",
    ], size=15)
    stat_card(s, Inches(7.2),  Inches(3.0), Inches(2.7), Inches(1.7),
              "73%", "of amateur players can't recall their last 5 drafts accurately", accent=ACCENT)
    stat_card(s, Inches(10.1), Inches(3.0), Inches(2.7), Inches(1.7),
              "2.4×", "improvement rate when prep is structured vs. ad-hoc", accent=PRIMARY)
    stat_card(s, Inches(7.2),  Inches(4.85), Inches(5.6), Inches(1.7),
              "1 roster", "is all it takes to start playing like an organized team", accent=SUCCESS)
    add_footer(s, 3, TOTAL)

    # 04 What the buyer gets
    s = base_slide(prs)
    add_slide_header(s, "03  ·  What you get",
                     "One roster. Eight workspaces. Zero spreadsheets.")
    items = [
        ("RST", "Roster",          "Players, roles, contact, status."),
        ("EVT", "Events",          "Tournaments, scrims, VOD reviews."),
        ("MTH", "Matches",         "Full results with maps and modes."),
        ("OPP", "Opponents",       "Linked head-to-head history."),
        ("HRO", "Heroes / Agents", "Pick-rates, win-rates, comps."),
        ("MAP", "Maps",            "Side, mode and map performance."),
        ("DFT", "Draft / B&P",     "What you ban, pick and beat."),
        ("STA", "Stats",           "Personal & team trends over time."),
    ]
    for i, (icon, title, body) in enumerate(items):
        col = i % 4
        row = i // 4
        x = Inches(0.6 + col * 3.1)
        y = Inches(3.0 + row * 2.0)
        feature_card(s, x, y, Inches(2.95), Inches(1.85), icon, title, body)
    add_footer(s, 4, TOTAL)

    # 05 Features included
    s = base_slide(prs)
    add_slide_header(s, "04  ·  Features included",
                     "Everything below ships with a single roster — no add-ons, no upsells.")
    blocks = [
        ("Roster & attendance", "Add players, set roles, track availability for every event."),
        ("Match logging",        "Capture maps, modes, score, result and notes in under 60 seconds."),
        ("Opponent profiles",    "Link every match to an opponent and see your true H2H history."),
        ("Hero / agent stats",   "Pick-rate, win-rate and synergy across your last N matches."),
        ("Map analytics",        "Side win-rate, attack vs. defense, mode-by-mode performance."),
        ("Draft / B&P insights", "What gets banned, what wins, what loses — over time."),
        ("Player comparisons",   "Stack two players head-to-head on any metric."),
        ("Filters & timeframes", "Slice by event, opponent, patch, map or window."),
        ("Dark / light themes",  "Built for late-night prep sessions and bright cafés alike."),
    ]
    for i, (t, b) in enumerate(blocks):
        col = i % 3
        row = i // 3
        x = Inches(0.6 + col * 4.15)
        y = Inches(3.0 + row * 1.35)
        add_rect(s, x, y, Inches(4.0), Inches(1.2), fill=SURFACE, line=LINE)
        add_rect(s, x, y, Inches(0.06), Inches(1.2), fill=PRIMARY)
        add_text(s, x + Inches(0.25), y + Inches(0.15), Inches(3.7), Inches(0.4),
                 t, size=12, bold=True, color=TEXT)
        add_text(s, x + Inches(0.25), y + Inches(0.55), Inches(3.7), Inches(0.6),
                 b, size=10, color=TEXT_MUTED, line_spacing=1.3)
    add_footer(s, 5, TOTAL)

    # 06 How it helps performance / prep
    s = base_slide(prs)
    add_slide_header(s, "05  ·  Performance & prep",
                     "From 'we'll figure it out' to 'we've prepared for this'.")
    bar_chart(s, Inches(0.6), Inches(3.0), Inches(7.5), Inches(3.7),
              "Match win-rate by preparation depth (illustrative)",
              [
                  ("No prep",         42, DANGER),
                  ("Quick scout",     54, WARN),
                  ("Structured prep", 67, PRIMARY),
                  ("Platform prep",   78, SUCCESS),
              ], max_val=100)
    add_bullets(s, Inches(8.4), Inches(3.0), Inches(4.5), Inches(4), [
        "Walk into matches with a plan, not a vibe",
        "Identify your weakest map and fix it",
        "Spot the comps your opponent loses to",
        "Catch your own bad habits before they cost you",
        "Build a season narrative you can actually review",
    ], size=12, gap_pt=4)
    add_footer(s, 6, TOTAL)

    # 07 How it looks / works (screenshot mockup)
    s = base_slide(prs)
    add_slide_header(s, "06  ·  How it looks",
                     "Built like a pro tool. Designed to be readable in a war room.")
    mockup_frame(s, Inches(0.6), Inches(2.95), Inches(8.0), Inches(4.0),
                 caption="Dashboard  ·  live product view",
                 image="Dashboard_users_1771832454665.png")
    add_bullets(s, Inches(8.9), Inches(3.0), Inches(4.0), Inches(4), [
        "Dark, premium UI",
        "Information-dense without being noisy",
        "Keyboard-friendly entry",
        "Fast filters across every view",
        "Mobile-friendly for live events",
    ], size=12)
    add_footer(s, 7, TOTAL)

    # 08 How it works (pipeline)
    s = base_slide(prs)
    add_slide_header(s, "07  ·  How it works",
                     "Four steps from sign-up to your first informed match.")
    diagram_pipeline(s, Inches(0.6), Inches(3.0), Inches(12.1), Inches(3.5), [
        ("01", "Create roster",   "Add your players, set roles and your competitive game."),
        ("02", "Log matches",     "Drop in results, maps, drafts and notes as you play."),
        ("03", "Review insights", "Trends, comps and matchup data appear automatically."),
        ("04", "Prep & repeat",   "Use the data to plan your next match and tighten the loop."),
    ])
    add_footer(s, 8, TOTAL)

    # 09 Example use case
    s = base_slide(prs)
    add_slide_header(s, "08  ·  Example use case",
                     "A solo competitive player preparing for an online cup.",
                     "Two weeks of prep. One roster. A measurable outcome.")
    timeline = [
        ("Day 1",  "Sets up roster, imports last 10 scrim results."),
        ("Day 3",  "Identifies a 31% win-rate on Mode X — flags it for practice."),
        ("Day 7",  "Scouts cup opponents — spots a recurring ban pattern."),
        ("Day 10", "Drafts three counter-comps in the planner."),
        ("Cup day", "Wins 4 / 5 best-of-threes. 2 of them on Mode X."),
    ]
    for i, (when, what) in enumerate(timeline):
        y = Inches(3.0 + i * 0.75)
        add_rect(s, Inches(0.6),  y, Inches(1.5), Inches(0.6), fill=SURFACE_2, line=LINE)
        add_text(s, Inches(0.6),  y, Inches(1.5), Inches(0.6),
                 when, size=12, bold=True, color=PRIMARY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(2.2),  y, Inches(10.5), Inches(0.6), fill=SURFACE, line=LINE)
        add_text(s, Inches(2.4),  y, Inches(10.2), Inches(0.6),
                 what, size=12, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(s, 9, TOTAL)

    # 10 Pricing framing
    s = base_slide(prs)
    add_slide_header(s, "09  ·  Pricing",
                     "One price per roster. Pick the plan that matches your game.")
    plans = [
        ("Lite",     "Free Fire, Fortnite, Apex, PUBG, EAFC, Rocket League, eFootball, Deadlock, TrackMania, fighting games",
         "$8 USD / month per roster",
         ["3 months: $18  (save 25%)",
          "6 months: $31  (save 35%)",
          "All features included",
          "Cancel anytime"], False),
        ("Standard", "Dota 2, CS:GO, LoL, Mobile Legends, Honor of Kings, TFT, CrossFire, CoD Mobile, The Finals, Brawl Stars",
         "$20 USD / month per roster",
         ["3 months: $45  (save 25%)",
          "6 months: $78  (save 35%)",
          "All features included",
          "Cancel anytime"], True),
        ("Premium",  "Valorant, Overwatch, Rainbow Six Siege, Call of Duty, Marvel Rivals",
         "$40 USD / month per roster",
         ["3 months: $90  (save 25%)",
          "6 months: $156 (save 35%)",
          "All features included",
          "Cancel anytime"], False),
    ]
    for i, (name, sub, price, feats, hi) in enumerate(plans):
        x = Inches(0.6 + i * 4.15)
        h = Inches(4.0)
        add_rect(s, x, Inches(3.0), Inches(4.0), h,
                 fill=SURFACE_2 if hi else SURFACE,
                 line=PRIMARY if hi else LINE,
                 line_w=1.5 if hi else 0.75)
        if hi:
            add_rect(s, x, Inches(2.78), Inches(1.4), Inches(0.3),
                     fill=PRIMARY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            add_text(s, x, Inches(2.78), Inches(1.4), Inches(0.3),
                     "MOST POPULAR", size=8, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.3), Inches(3.15), Inches(3.4), Inches(0.45),
                 name, size=16, bold=True, color=TEXT)
        add_text(s, x + Inches(0.3), Inches(3.6), Inches(3.4), Inches(0.7),
                 sub, size=9, color=TEXT_MUTED, line_spacing=1.2)
        add_text(s, x + Inches(0.3), Inches(4.35), Inches(3.4), Inches(0.5),
                 price, size=18, bold=True, color=PRIMARY if hi else TEXT)
        add_bullets(s, x + Inches(0.3), Inches(4.95), Inches(3.4), Inches(2.0),
                    feats, size=10, gap_pt=2)
    add_text(s, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.3),
             "Website-only access (no roster analytics): $40 USD / month flat.",
             size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_footer(s, 10, TOTAL)

    # 11 Manual vs Platform
    s = base_slide(prs)
    add_slide_header(s, "10  ·  Manual vs. Platform",
                     "What changes when you stop using spreadsheets and notes apps.")
    comparison_table(s, Inches(0.6), Inches(3.0), Inches(12.1), Inches(3.7),
                     ["Capability", "Spreadsheets / Notes", "Single Roster Plan"],
                     [
                         ["Match history",       "Manual entry, often forgotten", "Logged in under a minute, searchable"],
                         ["Opponent scouting",   "Scrolling old chats",           "Linked H2H profiles auto-built"],
                         ["Hero / map insights", "Calculated by hand (rarely)",   "Win-rate, pick-rate ready instantly"],
                         ["Draft analysis",      "Effectively impossible",        "Bans, picks and counters over time"],
                         ["Sharing with coach",  "Screenshots in Discord",        "Clean read-only view (with add-on)"],
                     ], highlight_col=2)
    add_footer(s, 11, TOTAL)

    # 12 FAQ / objections
    s = base_slide(prs)
    add_slide_header(s, "11  ·  FAQ & objections",
                     "The honest answers to what most buyers ask first.")
    faqs = [
        ("Is one roster really enough for me?",
         "Yes — if you only play one game competitively with one line-up, one roster covers everything."),
        ("Do I need to be on a team?",
         "No. Many users are solo competitive players using the roster as their personal record."),
        ("Will it work for my game?",
         "If your game has matches, maps/modes and characters/heroes, the templates already fit."),
        ("Can I export my data later?",
         "Yes. Your data stays yours; export is supported."),
        ("What if I outgrow it?",
         "You can upgrade to the Organizations plan and bring all your data with you."),
        ("How long until I see value?",
         "Most users get useful insights after logging 8–10 matches — usually within the first week."),
    ]
    for i, (q, a) in enumerate(faqs):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 6.2)
        y = Inches(3.0 + row * 1.35)
        add_rect(s, x, y, Inches(6.0), Inches(1.2), fill=SURFACE, line=LINE)
        add_text(s, x + Inches(0.25), y + Inches(0.15), Inches(5.7), Inches(0.4),
                 q, size=11, bold=True, color=PRIMARY)
        add_text(s, x + Inches(0.25), y + Inches(0.55), Inches(5.7), Inches(0.6),
                 a, size=10, color=TEXT_MUTED, line_spacing=1.3)
    add_footer(s, 12, TOTAL)

    # 13 Why now
    s = base_slide(prs)
    add_slide_header(s, "12  ·  Why now",
                     "Every match you play unprepared is data you'll never get back.")
    add_bullets(s, Inches(0.6), Inches(3.0), Inches(7.5), Inches(4),
                [
                    "Your competition is already tracking — quietly",
                    "Patches change the meta faster than memory can keep up",
                    "Off-season is the cheapest time to build the habit",
                    "You don't need a team to play like one",
                ], size=15)
    stat_card(s, Inches(8.7), Inches(3.0), Inches(4.0), Inches(1.7),
              "< 1 min", "to log a match in the platform", accent=PRIMARY)
    stat_card(s, Inches(8.7), Inches(4.85), Inches(4.0), Inches(1.7),
              "8–10 games", "until insights become actionable", accent=ACCENT)
    add_footer(s, 13, TOTAL)

    # 14 Risk reversal
    s = base_slide(prs)
    add_slide_header(s, "13  ·  Zero-risk start",
                     "We make it easy to try, easy to stay, easy to leave.")
    items = [
        ("01", "Cancel anytime",   "No long-term contracts or hidden fees."),
        ("02", "Your data is yours", "Export anytime in standard formats."),
        ("03", "Migrate up later", "Move to Organizations plan with no data loss."),
        ("04", "Live support",     "Real humans answer real questions."),
    ]
    for i, (icon, t, b) in enumerate(items):
        x = Inches(0.6 + i * 3.1)
        feature_card(s, x, Inches(3.0), Inches(2.95), Inches(3.5), icon, t, b)
    add_footer(s, 14, TOTAL)

    # 15 CTA
    cta_slide(
        prs,
        eyebrow="Take the next step",
        headline="Buy one roster.\nStart playing like a team of one.",
        subhead="Sign up, set up your roster in under five minutes,\nand log your first match tonight.",
        primary_label="START FREE TRIAL",
        secondary_label="BOOK A 15-MIN DEMO",
        footnote="Replace the buttons with your real URLs in Canva.",
    )

    # 16 Closing
    closing_slide(
        prs,
        headline="Talent wins games.\nData wins seasons.",
        line2="Roster Intelligence  ·  Single Roster Edition",
        contact="hello@yourdomain.com   ·   yourdomain.com   ·   @yourhandle",
    )

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)


# =====================================================================
#                       ORGANIZATIONS DECK (19 slides)
# =====================================================================
def build_organizations(out_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    TOTAL = 19

    # 01 Cover
    cover_slide(
        prs,
        eyebrow="Roster Intelligence  ·  Organization Edition",
        title="Run your esports org\nlike a professional team.",
        subtitle="A unified roster, analytics and scouting platform —\nbuilt for organizations, coaches and managers.",
        edition_label="Organizations  ·  Multi-Roster Platform",
    )

    # 02 The org problem
    s = base_slide(prs)
    add_slide_header(s, "01  ·  The problem",
                     "Most esports orgs run on group chats and gut feel.",
                     "The data exists. It's just trapped in screenshots, DMs and a coach's head.")
    add_bullets(s, Inches(0.6), Inches(3.0), Inches(6.0), Inches(4),
                [
                    "Each roster tracks things differently — or not at all",
                    "Match history lives in chat, not in a system",
                    "Scouting is restarted from scratch every event",
                    "Coaching staff can't see what management sees",
                    "Performance reviews are vibes, not numbers",
                ], size=14)
    stat_card(s, Inches(7.0),  Inches(3.0), Inches(2.8), Inches(1.7),
              "< 20%", "of amateur orgs keep structured match data", accent=DANGER)
    stat_card(s, Inches(10.0), Inches(3.0), Inches(2.8), Inches(1.7),
              "5+ tools", "the average org juggles per roster", accent=ACCENT)
    stat_card(s, Inches(7.0),  Inches(4.85), Inches(5.8), Inches(1.7),
              "0", "single source of truth for performance — until now", accent=PRIMARY)
    add_footer(s, 2, TOTAL)

    # 03 Inefficiencies in current workflow
    s = base_slide(prs)
    add_slide_header(s, "02  ·  Today's workflow is leaking value",
                     "Where time, knowledge and competitive edge get lost.")
    leaks = [
        ("Information silos",   "Coaches, analysts and managers see different data — or none."),
        ("Manual aggregation",  "Hours per week spent rebuilding spreadsheets from scratch."),
        ("Lost institutional knowledge", "Players leave; their match history leaves with them."),
        ("Inconsistent stats",  "Each game tracked with its own ad-hoc template."),
        ("Slow scouting",       "Opponent prep restarts before every event."),
        ("No accountability loop", "No clear way to measure if interventions worked."),
    ]
    for i, (t, b) in enumerate(leaks):
        col = i % 3
        row = i // 3
        x = Inches(0.6 + col * 4.15)
        y = Inches(3.0 + row * 1.85)
        add_rect(s, x, y, Inches(4.0), Inches(1.7), fill=SURFACE, line=LINE)
        add_rect(s, x, y, Inches(4.0), Inches(0.06), fill=ACCENT)
        add_text(s, x + Inches(0.25), y + Inches(0.2), Inches(3.7), Inches(0.4),
                 t, size=12, bold=True, color=TEXT)
        add_text(s, x + Inches(0.25), y + Inches(0.7), Inches(3.7), Inches(0.95),
                 b, size=10, color=TEXT_MUTED, line_spacing=1.3)
    add_footer(s, 3, TOTAL)

    # 04 Section: Platform
    section_divider(prs, 1, "Platform", "One platform.\nEvery roster, every game, every match.")

    # 05 Platform overview
    s = base_slide(prs)
    add_slide_header(s, "03  ·  Platform overview",
                     "A single workspace for every roster your org runs.")
    pillars = [
        ("ROS", "Roster Management", "Multi-roster structure, players, roles, contracts, attendance."),
        ("ANL", "Analytics",         "Hero, map, mode, draft and player performance over time."),
        ("SCT", "Scouting",          "Opponent profiles auto-built from your match history."),
        ("ORG", "Org Controls",      "Roles, permissions, audit trail and access by team."),
    ]
    for i, (icon, t, b) in enumerate(pillars):
        x = Inches(0.6 + i * 3.1)
        feature_card(s, x, Inches(3.0), Inches(2.95), Inches(3.6), icon, t, b)
    add_footer(s, 5, TOTAL)

    # 06 Multi-roster structure
    s = base_slide(prs)
    add_slide_header(s, "04  ·  Multi-roster structure",
                     "Built around how real esports orgs are organized.")
    # diagram: org -> rosters -> players
    add_rect(s, Inches(5.4), Inches(3.0), Inches(2.5), Inches(0.7),
             fill=PRIMARY, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(s, Inches(5.4), Inches(3.0), Inches(2.5), Inches(0.7),
             "Organization", size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rosters = ["Main · Game A", "Academy · Game A", "Roster · Game B", "Roster · Game C"]
    for i, r in enumerate(rosters):
        x = Inches(0.6 + i * 3.1)
        add_rect(s, x, Inches(4.4), Inches(2.95), Inches(0.7),
                 fill=SURFACE_2, line=LINE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(s, x, Inches(4.4), Inches(2.95), Inches(0.7),
                 r, size=11, bold=True, color=TEXT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # connector
        add_rect(s, Inches(0.6 + i * 3.1) + Inches(1.475) - Inches(0.01),
                 Inches(3.7), Inches(0.02), Inches(0.7), fill=LINE)
        # players row
        add_rect(s, x, Inches(5.55), Inches(2.95), Inches(1.2),
                 fill=SURFACE, line=LINE)
        add_text(s, x + Inches(0.2), Inches(5.6), Inches(2.6), Inches(0.4),
                 "Players · Roles · Stats", size=10, bold=True, color=PRIMARY)
        add_text(s, x + Inches(0.2), Inches(5.95), Inches(2.6), Inches(0.75),
                 "Game-aware templates, contracts, attendance and per-player history.",
                 size=9, color=TEXT_MUTED, line_spacing=1.25)
    # horizontal connector
    add_rect(s, Inches(0.6) + Inches(1.475), Inches(3.7),
             Inches(3.1 * 3), Inches(0.02), fill=LINE)
    add_footer(s, 6, TOTAL)

    # 07 Game templates / stats structure
    s = base_slide(prs)
    add_slide_header(s, "05  ·  Game templates & stats structure",
                     "Standardized data, customized per game.",
                     "Every game gets a template that captures the right metrics out of the box.")
    games = [
        ("Hero shooters",    "Heroes, roles, maps, modes, comps, KDA"),
        ("Tactical FPS",     "Agents, sides, maps, rounds, eco, plant/defuse"),
        ("MOBAs",            "Heroes, lanes, picks, bans, objectives, KDA"),
        ("Battle royale",    "Squad, drop, placement, kills, rotations"),
        ("Sports / arcade",  "Score, mode, possession, custom KPIs"),
        ("Custom",           "Define your own template per roster"),
    ]
    for i, (t, b) in enumerate(games):
        col = i % 3
        row = i // 3
        x = Inches(0.6 + col * 4.15)
        y = Inches(3.0 + row * 1.85)
        add_rect(s, x, y, Inches(4.0), Inches(1.7), fill=SURFACE, line=LINE)
        add_rect(s, x, y, Inches(0.06), Inches(1.7), fill=PRIMARY)
        add_text(s, x + Inches(0.25), y + Inches(0.2), Inches(3.7), Inches(0.4),
                 t, size=12, bold=True, color=TEXT)
        add_text(s, x + Inches(0.25), y + Inches(0.7), Inches(3.7), Inches(0.95),
                 b, size=10, color=TEXT_MUTED, line_spacing=1.3)
    add_footer(s, 7, TOTAL)

    # 08 Analytics & scouting
    s = base_slide(prs)
    add_slide_header(s, "06  ·  Analytics & scouting",
                     "Turn every match into intelligence — automatically.")
    bar_chart(s, Inches(0.6), Inches(3.0), Inches(7.5), Inches(3.7),
              "Coverage by analysis dimension (illustrative)",
              [
                  ("Player performance", 95, PRIMARY),
                  ("Map / mode trends",  92, PRIMARY),
                  ("Hero / agent meta",  88, PRIMARY),
                  ("Draft / B&P",        85, ACCENT),
                  ("Opponent profiles",  90, ACCENT),
              ], max_val=100)
    add_bullets(s, Inches(8.4), Inches(3.0), Inches(4.5), Inches(4),
                [
                    "Player & team trends across patches",
                    "Map and mode win-rates by side",
                    "Hero / agent picks, bans and counters",
                    "Opponent H2H built automatically",
                    "Draft phase analysis with priorities",
                ], size=12, gap_pt=4)
    add_footer(s, 8, TOTAL)

    # 09 Opponent insights
    s = base_slide(prs)
    add_slide_header(s, "07  ·  Opponent insights",
                     "Walk into every match knowing what they do — and what they don't.")
    add_bullets(s, Inches(0.6), Inches(3.0), Inches(6.0), Inches(4),
                [
                    "Auto-built opponent profiles from your match history",
                    "Most-picked heroes / agents and ban tendencies",
                    "Map preferences and side win-rates",
                    "Patterns by stage, time of day, day of week",
                    "Notes from coaches and analysts on the same record",
                ], size=14)
    mockup_frame(s, Inches(7.0), Inches(2.95), Inches(5.8), Inches(4.0),
                 caption="Opponent profile  ·  live product view",
                 image="Opponents_page_1775904047714.png")
    add_footer(s, 9, TOTAL)

    # 10 Draft / Map / Hero analysis
    s = base_slide(prs)
    add_slide_header(s, "08  ·  Draft / Map / Hero analysis",
                     "Three of the highest-leverage signals in competitive play.")
    cards = [
        ("Draft / B&P",
         "What gets banned first, what your roster wins on, and the comps you should target — by opponent."),
        ("Map performance",
         "Side, mode and map win-rates. Identify your strongest pool and your weakest gap."),
        ("Hero / agent meta",
         "Personal pick-rates vs. real win-rates. Stop riding heroes that aren't actually working."),
    ]
    for i, (t, b) in enumerate(cards):
        x = Inches(0.6 + i * 4.15)
        add_rect(s, x, Inches(3.0), Inches(4.0), Inches(3.7), fill=SURFACE, line=LINE)
        add_rect(s, x, Inches(3.0), Inches(4.0), Inches(0.08),
                 fill=PRIMARY if i != 1 else ACCENT)
        add_text(s, x + Inches(0.3), Inches(3.25), Inches(3.4), Inches(0.5),
                 t, size=15, bold=True, color=TEXT)
        add_text(s, x + Inches(0.3), Inches(3.85), Inches(3.4), Inches(2.7),
                 b, size=12, color=TEXT_MUTED, line_spacing=1.4)
    add_footer(s, 10, TOTAL)

    # 11 Section: Operate
    section_divider(prs, 2, "Operate", "Run the org.\nManage rosters, staff and access.")

    # 12 Team management value
    s = base_slide(prs)
    add_slide_header(s, "09  ·  Team management",
                     "The everyday workflows your coaches and managers actually do.")
    items = [
        ("ATT", "Attendance",       "Per event, per player, with reasons and trends."),
        ("EVT", "Event planning",   "Tournaments, scrims, VOD reviews — all linked to results."),
        ("VOD", "VOD references",   "Attach links and notes directly to the match."),
        ("CHT", "Internal chat",    "Light context where work actually happens."),
        ("PLR", "Player profiles",  "Roles, status, contact, performance in one place."),
        ("RPT", "Performance reports", "Generate weekly / monthly read-outs in seconds."),
    ]
    for i, (icon, t, b) in enumerate(items):
        col = i % 3
        row = i // 3
        x = Inches(0.6 + col * 4.15)
        y = Inches(3.0 + row * 1.85)
        feature_card(s, x, y, Inches(4.0), Inches(1.7), icon, t, b)
    add_footer(s, 12, TOTAL)

    # 13 Scaling across staff / teams
    s = base_slide(prs)
    add_slide_header(s, "10  ·  Scaling across staff & teams",
                     "Your org can grow without your spreadsheets multiplying.")
    diagram_pipeline(s, Inches(0.6), Inches(3.0), Inches(12.1), Inches(3.5), [
        ("01", "Add a roster",       "New game, new line-up — set up in minutes from a template."),
        ("02", "Invite staff",       "Coaches, analysts, managers, with per-roster access."),
        ("03", "Standardize data",   "Same metrics, same definitions, every team."),
        ("04", "Roll up to org",     "Aggregate across rosters for leadership read-outs."),
    ])
    add_footer(s, 13, TOTAL)

    # 14 Security / access / org benefits
    s = base_slide(prs)
    add_slide_header(s, "11  ·  Security & access",
                     "Granular control over who sees what — by role, by roster.")
    items = [
        ("01", "Role-based access", "Owner, manager, coach, analyst, player — distinct permissions."),
        ("02", "Per-roster scope",  "Limit a coach to their roster; managers see across."),
        ("03", "Audit visibility",  "Track key changes, additions and removals."),
        ("04", "Data ownership",    "Your data, your export, your retention controls."),
    ]
    for i, (icon, t, b) in enumerate(items):
        x = Inches(0.6 + i * 3.1)
        feature_card(s, x, Inches(3.0), Inches(2.95), Inches(3.7), icon, t, b)
    add_footer(s, 14, TOTAL)

    # 15 Why it matters competitively (comparison)
    s = base_slide(prs)
    add_slide_header(s, "12  ·  Why this matters competitively",
                     "What changes for the org once the platform is in place.")
    comparison_table(s, Inches(0.6), Inches(3.0), Inches(12.1), Inches(3.7),
                     ["Capability", "Today (manual)", "With the Platform"],
                     [
                         ["Cross-roster visibility", "None — every roster is its own island", "Unified org view across all rosters"],
                         ["Match history",           "Lost when players leave",                "Persistent, searchable, owned by the org"],
                         ["Scouting",                "Restarted every event",                  "Auto-built from your own match history"],
                         ["Performance review",     "Anecdotal, vibes-based",                  "Quantified, consistent, comparable"],
                         ["Onboarding new staff",    "Months of context-building",             "Hours — the data is already there"],
                     ], highlight_col=2)
    add_footer(s, 15, TOTAL)

    # 16 ROI framing
    s = base_slide(prs)
    add_slide_header(s, "13  ·  ROI framing",
                     "What organizations consistently get back.")
    stats = [
        ("8–12h",  "of analyst / coach time saved per roster, per month",   PRIMARY),
        ("100%",   "of match history retained when players move on",        SUCCESS),
        ("3×",     "faster opponent prep vs. manual scouting",              ACCENT),
        ("1",      "single source of truth for the entire org",             PRIMARY),
    ]
    for i, (v, l, c) in enumerate(stats):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 6.2)
        y = Inches(3.0 + row * 1.85)
        stat_card(s, x, y, Inches(6.0), Inches(1.7), v, l, accent=c)
    add_footer(s, 16, TOTAL)

    # 17 Pricing / adoption framing
    s = base_slide(prs)
    add_slide_header(s, "14  ·  Pricing",
                     "Per-roster pricing tiered by game. Scales linearly across your org.")
    plans = [
        ("Lite",     "Free Fire, Fortnite, Apex, PUBG, EAFC, Rocket League, eFootball, Deadlock, TrackMania, fighting games",
         "$8 USD / month per roster",
         ["3 months: $18  per roster (save 25%)",
          "6 months: $31  per roster (save 35%)",
          "Unlimited staff seats per roster",
          "All org features included"], False),
        ("Standard", "Dota 2, CS:GO, LoL, Mobile Legends, Honor of Kings, TFT, CrossFire, CoD Mobile, The Finals, Brawl Stars",
         "$20 USD / month per roster",
         ["3 months: $45  per roster (save 25%)",
          "6 months: $78  per roster (save 35%)",
          "Unlimited staff seats per roster",
          "All org features included"], True),
        ("Premium",  "Valorant, Overwatch, Rainbow Six Siege, Call of Duty, Marvel Rivals",
         "$40 USD / month per roster",
         ["3 months: $90  per roster (save 25%)",
          "6 months: $156 per roster (save 35%)",
          "Unlimited staff seats per roster",
          "All org features included"], False),
    ]
    for i, (name, sub, price, feats, hi) in enumerate(plans):
        x = Inches(0.6 + i * 4.15)
        h = Inches(4.0)
        add_rect(s, x, Inches(3.0), Inches(4.0), h,
                 fill=SURFACE_2 if hi else SURFACE,
                 line=PRIMARY if hi else LINE,
                 line_w=1.5 if hi else 0.75)
        if hi:
            add_rect(s, x, Inches(2.78), Inches(1.4), Inches(0.3),
                     fill=PRIMARY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            add_text(s, x, Inches(2.78), Inches(1.4), Inches(0.3),
                     "RECOMMENDED", size=8, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.3), Inches(3.15), Inches(3.4), Inches(0.45),
                 name, size=16, bold=True, color=TEXT)
        add_text(s, x + Inches(0.3), Inches(3.6), Inches(3.4), Inches(0.7),
                 sub, size=9, color=TEXT_MUTED, line_spacing=1.2)
        add_text(s, x + Inches(0.3), Inches(4.35), Inches(3.4), Inches(0.5),
                 price, size=18, bold=True, color=PRIMARY if hi else TEXT)
        add_bullets(s, x + Inches(0.3), Inches(4.95), Inches(3.4), Inches(2.0),
                    feats, size=10, gap_pt=2)
    add_text(s, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.3),
             "Website-only access (no roster analytics): $40 USD / month flat, org-wide.",
             size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_footer(s, 17, TOTAL)

    # 18 CTA
    cta_slide(
        prs,
        eyebrow="Move your org forward",
        headline="Bring every roster\ninto one professional system.",
        subhead="Pilot the platform with one roster. Roll it out across the org\nonce your staff has felt the difference.",
        primary_label="START AN ORG PILOT",
        secondary_label="BOOK A 30-MIN WALKTHROUGH",
        footnote="Replace the buttons with your real URLs in Canva.",
    )

    # 19 Closing
    closing_slide(
        prs,
        headline="Built for orgs\nthat play to win.",
        line2="Roster Intelligence  ·  Organization Edition",
        contact="hello@yourdomain.com   ·   yourdomain.com   ·   @yourhandle",
    )

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)


if __name__ == "__main__":
    build_individuals("deliverables/sales_decks/Individuals_Deck.pptx")
    build_organizations("deliverables/sales_decks/Organizations_Deck.pptx")
    print("Done.")
